#!/usr/bin/env python3
"""Build Chicks defense presentation — academic style, screenshot placeholders."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

TEMPLATE = "/Users/natala/Downloads/защита_пример.pptx"
OUTPUT_PATHS = [
    "/Users/natala/Downloads/защита_Chicks.pptx",
    "/Users/natala/Desktop/защита_Chicks.pptx",
]

FONT = "Calibri"
GRAY_FILL = RGBColor(0xEE, 0xEE, 0xEE)
GRAY_LINE = RGBColor(0xAA, 0xAA, 0xAA)
GRAY_TEXT = RGBColor(0x77, 0x77, 0x77)
CAPTION_COLOR = RGBColor(0x33, 0x33, 0x33)


def clear_slide_shapes(slide, keep_placeholders=False):
    """Remove all non-placeholder shapes added previously."""
    to_remove = []
    for shape in slide.shapes:
        if keep_placeholders and shape.is_placeholder:
            continue
        if not keep_placeholders:
            to_remove.append(shape)
        elif not shape.is_placeholder:
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape.element
        sp.getparent().remove(sp)


def set_title(slide, text, size=28):
  idx = 0
  for i, sh in enumerate(slide.shapes):
    if sh.is_placeholder and sh.placeholder_format.idx == 0:
      idx = i
      break
  shape = slide.shapes[idx]
  tf = shape.text_frame
  tf.clear()
  p = tf.paragraphs[0]
  p.text = text
  p.font.size = Pt(size)
  p.font.name = FONT
  p.font.bold = True


def add_text_block(slide, left, top, width, height, lines, base_size=13, bold_first=False):
  """lines: list of (text, size, bold, space_after) or plain strings."""
  box = slide.shapes.add_textbox(left, top, width, height)
  tf = box.text_frame
  tf.word_wrap = True
  tf.vertical_anchor = MSO_ANCHOR.TOP

  for i, item in enumerate(lines):
    if isinstance(item, str):
      text, size, bold, space = item, base_size, False, Pt(4)
    else:
      text, size, bold, space = item

    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.name = FONT
    p.font.bold = bold
    p.space_after = space
    p.line_spacing = 1.15
  return box


def add_screenshot_slot(slide, left, top, width, height, caption):
  rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
  rect.fill.solid()
  rect.fill.fore_color.rgb = GRAY_FILL
  rect.line.color.rgb = GRAY_LINE
  rect.line.width = Pt(1)

  tf = rect.text_frame
  tf.word_wrap = True
  tf.vertical_anchor = MSO_ANCHOR.MIDDLE
  p = tf.paragraphs[0]
  p.alignment = PP_ALIGN.CENTER
  p.text = "Вставить скриншот"
  p.font.size = Pt(11)
  p.font.color.rgb = GRAY_TEXT
  p.font.name = FONT

  cap_h = Inches(0.32)
  cap = slide.shapes.add_textbox(left, top + height + Inches(0.04), width, cap_h)
  cp = cap.text_frame.paragraphs[0]
  cp.text = caption
  cp.font.size = Pt(9)
  cp.font.italic = True
  cp.font.name = FONT
  cp.font.color.rgb = CAPTION_COLOR
  cp.alignment = PP_ALIGN.CENTER


def add_grid(slide, items, left, top, cols, cell_w, cell_h, gap_x, gap_y):
  """items: list of caption strings."""
  for i, caption in enumerate(items):
    row, col = divmod(i, cols)
    x = left + col * (cell_w + gap_x)
    y = top + row * (cell_h + gap_y + Inches(0.36))
    add_screenshot_slot(slide, x, y, cell_w, cell_h, caption)


def prep_slide(slide, title=None):
  """Clear body placeholders content; set title if given."""
  for sh in slide.shapes:
    if sh.is_placeholder and sh.has_text_frame:
      idx = sh.placeholder_format.idx
      if idx == 0 and title:
        tf = sh.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.name = FONT
        p.font.bold = True
      elif idx != 0:
        tf = sh.text_frame
        tf.clear()


def build():
  prs = Presentation(TEMPLATE)
  layout = prs.slide_layouts[1]  # Заголовок и объект

  # ── Слайд 1. Титульный ──────────────────────────────────────────────
  s1 = prs.slides[0]
  prep_slide(s1)
  for sh in s1.shapes:
    if sh.is_placeholder and sh.placeholder_format.idx == 0:
      tf = sh.text_frame
      tf.clear()
      p = tf.paragraphs[0]
      p.text = "Тема: Разработка мобильного приложения персонального ИИ-стилиста «Chicks»"
      p.font.size = Pt(24)
      p.font.name = FONT
      p.font.bold = True
    if sh.is_placeholder and sh.placeholder_format.idx == 1:
      tf = sh.text_frame
      tf.clear()
      p = tf.paragraphs[0]
      p.text = "Выполнила: ФИО\nГруппа: ___\n\nНаучный руководитель:\nЭмирова З.М."
      p.font.size = Pt(18)
      p.font.name = FONT

  add_screenshot_slot(
    s1, Inches(6.8), Inches(2.6), Inches(2.8), Inches(2.4),
    "Рис. 1. Мобильное приложение «Chicks» — персональный ИИ-стилист",
  )

  # ── Слайд 2. Актуальность ───────────────────────────────────────────
  s2 = prs.slides[1]
  prep_slide(s2, "АКТУАЛЬНОСТЬ")
  add_text_block(
    s2, Inches(0.5), Inches(1.1), Inches(5.8), Inches(2.8),
    [
      ("Ежедневный выбор одежды отнимает время и вызывает стресс у большинства пользователей.", 13, False, Pt(6)),
      ("Услуги профессионального стилиста дороги и недоступны широкой аудитории.", 13, False, Pt(6)),
      ("Покупки без учёта имеющегося гардероба приводят к неиспользуемым вещам.", 13, False, Pt(6)),
      ("Развитие технологий искусственного интеллекта позволяет создавать персональные fashion-рекомендации в мобильном приложении.", 13, False, Pt(6)),
      ("Растёт спрос на приложения с учётом типа фигуры, цветотипа и личного гардероба.", 13, False, Pt(10)),
      ("Вывод: разработка приложения, объединяющего цифровой гардероб, анализ стиля и ИИ-консультанта, является актуальной.", 13, True, Pt(4)),
    ],
  )
  add_grid(
    s2,
    [
      "Рис. 2. ИИ-стилист помогает выбрать образ",
      "Рис. 3. Цифровой гардероб пользователя",
    ],
    Inches(0.5), Inches(3.5), 2, Inches(4.3), Inches(1.55), Inches(0.35), Inches(0.4),
  )

  # ── Слайд 3. Цель и задачи ───────────────────────────────────────────
  s3 = prs.slides[2]
  prep_slide(s3)
  add_text_block(
    s3, Inches(0.45), Inches(0.35), Inches(5.9), Inches(5.0),
    [
      ("ЦЕЛЬ", 15, True, Pt(4)),
      ("Разработать кроссплатформенное мобильное приложение «Chicks» для персональных рекомендаций по стилю на основе цифрового гардероба и технологий искусственного интеллекта.", 12, False, Pt(10)),
      ("ЗАДАЧИ", 15, True, Pt(4)),
      ("1. Проанализировать предметную область и существующие fashion-приложения.", 12, False, Pt(3)),
      ("2. Сформировать функциональные и нефункциональные требования.", 12, False, Pt(3)),
      ("3. Спроектировать архитектуру мобильного приложения.", 12, False, Pt(3)),
      ("4. Реализовать модуль гардероба с распознаванием одежды по фотографии.", 12, False, Pt(3)),
      ("5. Реализовать чат с ИИ-стилистом с учётом гардероба, погоды и контекста.", 12, False, Pt(3)),
      ("6. Реализовать анализ стиля, онбординг и мультиязычность.", 12, False, Pt(3)),
      ("7. Провести тестирование и оценить результаты разработки.", 12, False, Pt(10)),
      ("ОБЪЕКТ исследования — процесс создания мобильных приложений в сфере fashion-tech.", 12, False, Pt(4)),
      ("ПРЕДМЕТ исследования — методы разработки ИИ-приложения для подбора образов и управления гардеробом.", 12, False, Pt(4)),
    ],
  )
  add_screenshot_slot(
    s3, Inches(6.6), Inches(0.8), Inches(3.0), Inches(4.6),
    "Рис. 4. Основные модули приложения на главном экране",
  )

  # ── Слайд 4. Глава 1 ─────────────────────────────────────────────────
  s4 = prs.slides[3]
  prep_slide(s4, "ГЛАВА 1. Теоретические основы")
  add_text_block(
    s4, Inches(0.45), Inches(1.05), Inches(5.5), Inches(4.2),
    [
      ("1.1. Понятие персонального стайлинга и цифрового гардероба.", 13, False, Pt(8)),
      ("1.2. Обзор аналогов: Pinterest, Stylebook, Whering, ChatGPT. Отличие «Chicks»: объединение гардероба, ИИ-чата, анализа стиля и учёта цветотипа в одном приложении.", 13, False, Pt(8)),
      ("1.3. Технологии: Flutter (интерфейс), Firebase (авторизация, база данных, фото), OpenAI (чат и распознавание одежды), BLoC и Hive (логика и локальное хранение).", 12, False, Pt(8)),
      ("1.4. Требования к системе: функциональные (гардероб, чат, профиль, избранное, история) и нефункциональные (производительность, безопасность, локализация).", 13, False, Pt(4)),
    ],
  )
  add_grid(
    s4,
    [
      "Рис. 5. Определение цветотипа пользователя",
      "Рис. 6. Определение типа фигуры",
      "Рис. 7. Параметры стиля в профиле",
    ],
    Inches(6.2), Inches(1.05), 1, Inches(3.5), Inches(1.35), Inches(0), Inches(0.45),
  )

  # ── Слайд 5. Глава 2 ─────────────────────────────────────────────────
  s5 = prs.slides[4]
  prep_slide(s5, "ГЛАВА 2. Практическая реализация")
  add_text_block(
    s5, Inches(0.45), Inches(1.05), Inches(9.2), Inches(1.55),
    [
      ("2.1. Архитектура приложения построена по модульному принципу с разделением интерфейса, бизнес-логики и данных.", 12, False, Pt(5)),
      ("2.2. Реализованы модули: гардероб, ИИ-чат, анализ стиля, профиль, избранное и история, локализация на трёх языках.", 12, False, Pt(5)),
      ("2.3. Стек технологий: Flutter, Dart, OpenAI, Firebase. Результат: рабочее приложение для Android и iOS.", 12, False, Pt(4)),
    ],
  )
  add_grid(
    s5,
    [
      "Рис. 8. Экран цифрового гардероба",
      "Рис. 9. Распознавание одежды по фотографии",
      "Рис. 10. Детальная карточка элемента гардероба",
      "Рис. 11. Подбор образа из гардероба",
    ],
    Inches(0.45), Inches(2.55), 2, Inches(4.4), Inches(1.35), Inches(0.3), Inches(0.38),
  )

  # ── Слайд 6 → Глава 2 продолжение ────────────────────────────────────
  s6 = prs.slides[5]
  prep_slide(s6, "ГЛАВА 2 (продолжение)")
  add_grid(
    s6,
    [
      "Рис. 12. Персональные инсайты и учёт погоды",
      "Рис. 13. Анализ гардероба",
    ],
    Inches(1.2), Inches(1.5), 2, Inches(3.8), Inches(3.0), Inches(0.5), Inches(0.5),
  )

  # ── Слайд 7. Демонстрация ────────────────────────────────────────────
  s7 = prs.slides.add_slide(layout)
  prep_slide(s7, "ДЕМОНСТРАЦИЯ ИНТЕРФЕЙСА")
  add_grid(
    s7,
    [
      "Рис. 15. Авторизация пользователя",
      "Рис. 16. Контекстные подсказки для стилиста",
      "Рис. 17. Фильтрация элементов гардероба",
      "Рис. 18. Профиль пользователя",
    ],
    Inches(0.5), Inches(1.15), 2, Inches(4.4), Inches(1.75), Inches(0.35), Inches(0.45),
  )

  # ── Слайд 8. Заключение ──────────────────────────────────────────────
  s8 = prs.slides.add_slide(layout)
  prep_slide(s8, "ЗАКЛЮЧЕНИЕ")
  add_text_block(
    s8, Inches(0.45), Inches(1.05), Inches(5.6), Inches(4.3),
    [
      ("В ходе работы проанализирована предметная область fashion-tech и обоснована актуальность темы.", 12, False, Pt(5)),
      ("Спроектирована и реализована архитектура мобильного приложения «Chicks».", 12, False, Pt(5)),
      ("Создан модуль цифрового гардероба с распознаванием вещей по фотографии.", 12, False, Pt(5)),
      ("Реализован интеллектуальный чат-стилист с учётом гардероба, погоды и профиля пользователя.", 12, False, Pt(5)),
      ("Добавлены анализ стиля, онбординг, избранное, история образов и поддержка трёх языков.", 12, False, Pt(5)),
      ("Проведено тестирование на платформах Android и iOS.", 12, False, Pt(10)),
      ("Практическая значимость: приложение помогает пользователям быстрее выбирать образы и рациональнее использовать гардероб.", 12, False, Pt(8)),
      ("Перспективы развития: подписка, интеграция с магазинами, push-уведомления.", 12, False, Pt(10)),
      ("Спасибо за внимание!", 14, True, Pt(4)),
    ],
  )
  add_screenshot_slot(
    s8, Inches(6.5), Inches(1.05), Inches(3.1), Inches(4.2),
    "Рис. 14. Итоговый вид приложения «Chicks»",
  )

  for path in OUTPUT_PATHS:
    prs.save(path)
    print(f"Saved: {path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
  build()
